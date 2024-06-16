from pptx import Presentation
from gtts import gTTS
import os

# Crear la carpeta 'audio' si no existe
audio_dir = 'audio'
os.makedirs(audio_dir, exist_ok=True)

# Cargar la presentación de PowerPoint
ppt_path = 'Proyecto_Web_DAM.pptx'  # Cambia esto por la ruta de tu archivo
prs = Presentation(ppt_path)

# Texto para cada diapositiva
narrations = [
    "Hola, mi nombre es Mohamed ElDerkaoui Merzouk Bedoukha y esta es la presentación de mi proyecto sobre la mejora y ampliación de la página web de DAM Servicios Informáticos utilizando Wagtail CMS. Este proyecto se lleva a cabo desde el 1 de abril de 2024 hasta el 15 de junio de 2024.",
    "El objetivo principal de este proyecto es mejorar y ampliar la página web de DAM Servicios Informáticos utilizando Wagtail CMS, con el fin de proporcionar una experiencia en línea más completa y satisfactoria para nuestros clientes. Los objetivos específicos incluyen: Incorporar nuevos servicios, mejorar la usabilidad y actualizar el contenido.",
    "El alcance del proyecto incluye el diseño y desarrollo de nuevas páginas web, integración de herramientas interactivas, optimización del contenido existente y mejora de la navegación y usabilidad.",
    "Durante la fase de planificación de este proyecto, se llevaron a cabo análisis exhaustivos de diversas soluciones y alternativas para mejorar la página web de DAM Servicios Informáticos. Tras analizar las características, la flexibilidad y la escalabilidad de cada plataforma, se determinó que Wagtail CMS era la opción más adecuada.",
    "La planificación del proyecto se divide en varias fases: Análisis y requisitos (Semanas 1-2), Diseño (Semana 3), Desarrollo (Semanas 4-8), Pruebas (Semanas 9-10), Lanzamiento (Semana 11) y Post-lanzamiento (Monitoreo y mantenimiento inicial).",
    "El presupuesto estimado del proyecto incluye costos de diseño web, desarrollo web, alojamiento web y gastos generales, sumando un total de 7115.64 euros.",
    "Las áreas de gestión del proyecto incluyen gestión del alcance, gestión del tiempo, gestión del costo, gestión de la calidad, gestión de los recursos humanos, gestión de las comunicaciones y gestión de los riesgos.",
    "La aplicación de módulos incluye el desarrollo web con HTML5, CSS3 y JavaScript, desarrollo backend con Python y Django, gestión de bases de datos con PostgreSQL, y diseño y usabilidad aplicando principios de diseño para mejorar la experiencia del usuario.",
    "Referencias web utilizadas incluyen la documentación de Wagtail CMS y recursos educativos como lmorillas.github.io y wagtail.org.",
    "En conclusión, el proyecto ha logrado sus objetivos de mejorar y ampliar la página web de DAM Servicios Informáticos, proporcionando una experiencia en línea más completa y satisfactoria para nuestros clientes."
]

# Generar archivos de audio y añadirlos a las diapositivas
for i, slide in enumerate(prs.slides):
    if i < len(narrations):
        text = narrations[i]
        tts = gTTS(text=text, lang='es')
        audio_path = os.path.join(audio_dir, f"slide_{i+1}.mp3")
        tts.save(audio_path)

        # Añadir el archivo de audio a la diapositiva
        with open(audio_path, 'rb') as audio_file:
            slide.shapes.add_movie(audio_file, 0, 0, 0, 0, mime_type='audio/mp3')

# Guardar la presentación modificada
output_path = '/mnt/data/Modified_Proyecto_Web_DAM.pptx'  # Cambia esto por la ruta de salida
prs.save(output_path)

print(f"Presentación guardada en {output_path}")